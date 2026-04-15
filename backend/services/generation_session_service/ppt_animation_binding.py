from __future__ import annotations

from pathlib import Path
from typing import Any

from PIL import Image

try:
    from pptx import Presentation
except Exception:  # pragma: no cover - optional dependency guard
    Presentation = None

_BACKEND_ROOT = Path(__file__).resolve().parents[2]
_ANIMATION_MARKER_PREFIX = "spectra-animation:"


def _artifact_storage_path(artifact: Any) -> Path:
    raw_path = str(getattr(artifact, "storagePath", "") or "").strip()
    if not raw_path:
        raise FileNotFoundError("Artifact has no storagePath")

    candidate = Path(raw_path)
    if candidate.is_absolute():
        return candidate

    cwd_candidate = Path.cwd() / candidate
    if cwd_candidate.exists():
        return cwd_candidate

    backend_candidate = _BACKEND_ROOT / candidate
    return backend_candidate


def _slot_bounds(
    slot: str, slide_width: int, slide_height: int
) -> tuple[int, int, int, int]:
    margin_x = int(slide_width * 0.05)
    margin_y = int(slide_height * 0.06)

    if slot == "bottom-panel":
        box_width = int(slide_width * 0.72)
        box_height = int(slide_height * 0.22)
        left = max(margin_x, int((slide_width - box_width) / 2))
        top = slide_height - box_height - margin_y
        return left, top, box_width, box_height

    if slot == "right-panel":
        box_width = int(slide_width * 0.3)
        box_height = int(slide_height * 0.52)
        left = slide_width - box_width - margin_x
        top = max(margin_y, int((slide_height - box_height) / 2))
        return left, top, box_width, box_height

    box_width = int(slide_width * 0.28)
    box_height = int(slide_height * 0.24)
    left = slide_width - box_width - margin_x
    top = slide_height - box_height - margin_y
    return left, top, box_width, box_height


def _fit_image_into_box(
    image_width: int,
    image_height: int,
    box_left: int,
    box_top: int,
    box_width: int,
    box_height: int,
) -> tuple[int, int, int, int]:
    safe_width = max(1, int(image_width or 0))
    safe_height = max(1, int(image_height or 0))
    scale = min(box_width / safe_width, box_height / safe_height)
    target_width = max(1, int(safe_width * scale))
    target_height = max(1, int(safe_height * scale))
    left = box_left + max(0, int((box_width - target_width) / 2))
    top = box_top + max(0, int((box_height - target_height) / 2))
    return left, top, target_width, target_height


def _shape_marker(animation_artifact_id: str) -> str:
    return f"{_ANIMATION_MARKER_PREFIX}{animation_artifact_id}"


def _remove_existing_animation_shapes(slide, marker: str) -> None:
    removable_shapes = []
    for shape in slide.shapes:
        element = getattr(shape, "_element", None)
        if element is None:
            continue
        c_nv_pr = getattr(getattr(element, "nvPicPr", None), "cNvPr", None)
        descr = c_nv_pr.get("descr") if c_nv_pr is not None else ""
        if descr == marker or str(getattr(shape, "name", "") or "") == marker:
            removable_shapes.append(shape)

    for shape in removable_shapes:
        element = shape._element
        element.getparent().remove(element)


def _tag_picture_shape(picture, marker: str) -> None:
    try:
        picture.name = marker
    except Exception:  # pragma: no cover - best effort only
        pass

    c_nv_pr = getattr(getattr(picture._element, "nvPicPr", None), "cNvPr", None)
    if c_nv_pr is not None:
        c_nv_pr.set("name", marker)
        c_nv_pr.set("descr", marker)


def apply_animation_placement_to_ppt_artifact(
    *,
    animation_artifact: Any,
    ppt_artifact: Any,
    placement_records: list[dict[str, Any]],
) -> str:
    if Presentation is None:
        raise RuntimeError("python-pptx is required for PPT animation placement")

    ppt_path = _artifact_storage_path(ppt_artifact)
    gif_path = _artifact_storage_path(animation_artifact)
    if not ppt_path.exists():
        raise FileNotFoundError(f"PPT artifact file not found: {ppt_path}")
    if not gif_path.exists():
        raise FileNotFoundError(f"Animation artifact file not found: {gif_path}")

    with Image.open(gif_path) as image:
        image_width, image_height = image.size

    presentation = Presentation(str(ppt_path))
    marker = _shape_marker(str(getattr(animation_artifact, "id", "") or "").strip())

    for record in placement_records:
        page_number = int(record.get("page_number") or 0)
        if page_number < 1 or page_number > len(presentation.slides):
            continue
        slot = str(record.get("slot") or "bottom-right").strip().lower()
        slide = presentation.slides[page_number - 1]
        _remove_existing_animation_shapes(slide, marker)

        left, top, box_width, box_height = _slot_bounds(
            slot, presentation.slide_width, presentation.slide_height
        )
        left, top, width, height = _fit_image_into_box(
            image_width,
            image_height,
            left,
            top,
            box_width,
            box_height,
        )
        picture = slide.shapes.add_picture(
            str(gif_path),
            left=left,
            top=top,
            width=width,
            height=height,
        )
        _tag_picture_shape(picture, marker)

    presentation.save(str(ppt_path))
    return str(ppt_path)
