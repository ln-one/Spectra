"""
file_parser.py 回归测试。

验证 extract_text_for_rag() 在可插拔架构改造后行为与 MVP 完全一致。
测试文件在 tmp_path 中动态生成，不提交二进制 fixture。
"""

from pathlib import Path

from services.file_parser import extract_text_for_rag

# ---- 纯文本 ----


def test_extract_text_for_rag_txt_file(tmp_path: Path):
    file_path = tmp_path / "notes.txt"
    file_path.write_text("第一行内容\n第二行内容", encoding="utf-8")

    # 生产流程中 _resolve_file_type() 将 txt 归类为 "word"
    text, details = extract_text_for_rag(str(file_path), "notes.txt", "word")

    assert "第一行内容" in text
    assert details["text_length"] > 0


def test_extract_text_for_rag_md_file(tmp_path: Path):
    file_path = tmp_path / "readme.md"
    file_path.write_text("# 标题\n正文内容", encoding="utf-8")

    # 生产流程中 _resolve_file_type() 将 md 归类为 "word"
    text, details = extract_text_for_rag(str(file_path), "readme.md", "word")

    assert "标题" in text
    assert details["text_length"] > 0


def test_extract_text_for_rag_csv_file(tmp_path: Path):
    file_path = tmp_path / "data.csv"
    file_path.write_text("姓名,分数\n张三,90", encoding="utf-8")

    # 生产流程中 _resolve_file_type() 将 csv 归类为 "word"
    text, details = extract_text_for_rag(str(file_path), "data.csv", "word")

    assert "张三" in text
    assert details["text_length"] > 0


# ---- 图片/视频占位 ----


def test_extract_text_for_rag_image_placeholder(tmp_path: Path):
    file_path = tmp_path / "figure.png"
    file_path.write_bytes(b"\x89PNG\r\n")

    text, details = extract_text_for_rag(str(file_path), "figure.png", "image")

    assert "图片资料" in text
    assert "figure.png" in text
    assert details["images_extracted"] == 1


def test_extract_text_for_rag_video_placeholder(tmp_path: Path):
    file_path = tmp_path / "lecture.mp4"
    file_path.write_bytes(b"\x00\x00\x00\x1cftyp")

    text, details = extract_text_for_rag(str(file_path), "lecture.mp4", "video")

    assert "lecture.mp4" in text
    assert "capability_status" in details
    assert details["capability_status"]["capability"] == "video_understanding"
    assert details["text_length"] > 0


# ---- PDF 动态生成 ----


def test_extract_text_for_rag_pdf(tmp_path: Path):
    """动态生成最小 PDF 并验证解析。"""
    from pypdf import PdfWriter

    writer = PdfWriter()
    writer.add_blank_page(width=72, height=72)
    pdf_path = tmp_path / "test.pdf"
    with open(pdf_path, "wb") as f:
        writer.write(f)

    text, details = extract_text_for_rag(str(pdf_path), "test.pdf", "pdf")

    assert isinstance(text, str)
    assert details["pages_extracted"] == 1
    assert "text_length" in details


def test_extract_text_for_rag_pdf_corrupted(tmp_path: Path):
    """损坏 PDF 应返回空文本，不抛异常。"""
    pdf_path = tmp_path / "bad.pdf"
    pdf_path.write_bytes(b"not a pdf content at all")

    text, details = extract_text_for_rag(str(pdf_path), "bad.pdf", "pdf")

    assert text == ""
    assert details["pages_extracted"] == 0


# ---- DOCX 动态生成 ----


def test_extract_text_for_rag_docx(tmp_path: Path):
    """动态生成最小 DOCX 并验证解析。"""
    from docx import Document

    doc = Document()
    doc.add_paragraph("测试段落一")
    doc.add_paragraph("测试段落二")
    docx_path = tmp_path / "test.docx"
    doc.save(str(docx_path))

    text, details = extract_text_for_rag(str(docx_path), "test.docx", "word")

    assert "测试段落一" in text
    assert "测试段落二" in text
    assert details["pages_extracted"] >= 1
    assert details["text_length"] > 0


def test_extract_text_for_rag_docx_corrupted(tmp_path: Path):
    """损坏 DOCX 应返回空文本，不抛异常。"""
    docx_path = tmp_path / "bad.docx"
    docx_path.write_bytes(b"not a docx")

    text, details = extract_text_for_rag(str(docx_path), "bad.docx", "word")

    assert text == ""
    # 与 MVP 行为一致：DOCX 失败时仍返回 pages_extracted=1
    assert details["pages_extracted"] == 1


# ---- PPTX 动态生成 ----


def test_extract_text_for_rag_pptx(tmp_path: Path):
    """动态生成最小 PPTX 并验证解析。"""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    txBox.text_frame.text = "幻灯片文本内容"
    pptx_path = tmp_path / "test.pptx"
    prs.save(str(pptx_path))

    text, details = extract_text_for_rag(str(pptx_path), "test.pptx", "ppt")

    assert "幻灯片文本内容" in text
    assert details["pages_extracted"] == 1
    assert details["text_length"] > 0


def test_extract_text_for_rag_pptx_corrupted(tmp_path: Path):
    """损坏 PPTX 应返回空文本，不抛异常。"""
    pptx_path = tmp_path / "bad.pptx"
    pptx_path.write_bytes(b"not a pptx")

    text, details = extract_text_for_rag(str(pptx_path), "bad.pptx", "ppt")

    assert text == ""
    assert details["pages_extracted"] == 0


# ---- 接口签名兼容性 ----


def test_extract_text_for_rag_signature():
    """确保函数签名与 MVP 一致。"""
    import inspect

    sig = inspect.signature(extract_text_for_rag)
    params = list(sig.parameters.keys())
    assert params == ["filepath", "filename", "file_type"]


def test_extract_text_for_rag_fallback_to_local_when_provider_unsupported(
    tmp_path: Path, monkeypatch
):
    """当配置 provider 不支持 file_type 时，应自动回退到 local。"""
    import services.file_parser as file_parser_module
    from services.parsers.local_provider import LocalProvider

    class _UnsupportedProvider:
        name = "mineru"

        def supports(self, file_type: str) -> bool:
            return False

        def extract_text(self, filepath: str, filename: str, file_type: str):
            raise AssertionError("unsupported provider should not be used")

    calls = {"local": 0}

    def _fake_get_parser(provider_name=None):
        if provider_name == "local":
            calls["local"] += 1
            return LocalProvider()
        return _UnsupportedProvider()

    monkeypatch.setattr(file_parser_module, "get_parser", _fake_get_parser)

    # 使用 .docx 文件：生产 _resolve_file_type() 将 docx 归类为 "word"，
    # 且 .docx 不会被入口层纯文本短路，能真正走到 provider 回退逻辑。
    from docx import Document as DocxDocument

    docx_path = tmp_path / "notes.docx"
    _doc = DocxDocument()
    _doc.add_paragraph("fallback local content")
    _doc.save(str(docx_path))

    text, details = extract_text_for_rag(str(docx_path), "notes.docx", "word")

    assert "fallback local content" in text
    assert details["text_length"] > 0
    assert calls["local"] == 1


def test_extract_text_for_rag_fallback_to_llamaparse_when_mineru_unavailable(
    tmp_path: Path, monkeypatch
):
    import services.file_parser as file_parser_module

    class _FakeUnavailableProvider:
        name = "local"  # registry fallback behavior when provider unavailable

        def supports(self, _file_type: str) -> bool:
            return True

        def extract_text(self, _filepath: str, _filename: str, _file_type: str):
            return "", {"text_length": 0}

    class _FakeLlamaProvider:
        name = "llamaparse"

        def supports(self, _file_type: str) -> bool:
            return True

        def extract_text(self, _filepath: str, _filename: str, _file_type: str):
            return "llama parse content", {"pages_extracted": 1, "text_length": 19}

    calls = {"mineru": 0, "llamaparse": 0, "local": 0}

    def _fake_get_parser(provider_name=None):
        if provider_name == "mineru":
            calls["mineru"] += 1
            return _FakeUnavailableProvider()
        if provider_name == "llamaparse":
            calls["llamaparse"] += 1
            return _FakeLlamaProvider()
        calls["local"] += 1
        from services.parsers.local_provider import LocalProvider

        return LocalProvider()

    monkeypatch.setattr(file_parser_module, "get_parser", _fake_get_parser)
    monkeypatch.setenv("DOCUMENT_PARSER", "mineru")

    from docx import Document as DocxDocument

    docx_path = tmp_path / "notes.docx"
    _doc = DocxDocument()
    _doc.add_paragraph("fallback content")
    _doc.save(str(docx_path))

    text, details = extract_text_for_rag(str(docx_path), "notes.docx", "word")

    assert text == "llama parse content"
    assert details["capability_status"]["status"] == "degraded"
    assert details["capability_status"]["provider"] == "llamaparse"
    assert details["capability_status"]["fallback_target"] == "llamaparse"
    assert calls["mineru"] == 1
    assert calls["llamaparse"] == 1


def test_extract_text_for_rag_unknown_provider_falls_back_to_local(
    tmp_path: Path, monkeypatch
):
    import services.file_parser as file_parser_module
    from services.parsers.local_provider import LocalProvider

    class _ResolvedLocalProvider(LocalProvider):
        name = "local"

    calls = {"unknown": 0, "local": 0}

    def _fake_get_parser(provider_name=None):
        if provider_name == "unknown-provider":
            calls["unknown"] += 1
            return _ResolvedLocalProvider()
        if provider_name == "local":
            calls["local"] += 1
            return _ResolvedLocalProvider()
        return _ResolvedLocalProvider()

    monkeypatch.setattr(file_parser_module, "get_parser", _fake_get_parser)
    monkeypatch.setenv("DOCUMENT_PARSER", "unknown-provider")

    from docx import Document as DocxDocument

    docx_path = tmp_path / "notes.docx"
    _doc = DocxDocument()
    _doc.add_paragraph("local fallback content")
    _doc.save(str(docx_path))

    text, details = extract_text_for_rag(str(docx_path), "notes.docx", "word")

    assert "local fallback content" in text
    assert details["capability_status"]["status"] == "degraded"
    assert details["capability_status"]["provider"] == "local"
    assert details["capability_status"]["fallback_target"] == "local"
    assert calls["unknown"] == 1
    assert calls["local"] == 1


def test_extract_text_for_rag_auto_routes_pdf_to_mineru(tmp_path: Path, monkeypatch):
    import services.file_parser as file_parser_module

    class _FakeMineruProvider:
        name = "mineru"

        def supports(self, _file_type: str) -> bool:
            return True

        def extract_text(self, _filepath: str, _filename: str, _file_type: str):
            return "mineru parsed content", {"pages_extracted": 1, "text_length": 21}

    calls: list[str] = []

    def _fake_get_parser(provider_name=None):
        calls.append(str(provider_name))
        if provider_name == "mineru":
            return _FakeMineruProvider()
        raise AssertionError(f"unexpected provider request: {provider_name}")

    pdf_path = tmp_path / "sample.pdf"
    pdf_path.write_bytes(b"%PDF-1.4\n")

    monkeypatch.setenv("DOCUMENT_PARSER", "auto")
    monkeypatch.setattr(file_parser_module, "get_parser", _fake_get_parser)

    text, details = extract_text_for_rag(str(pdf_path), "sample.pdf", "pdf")

    assert text == "mineru parsed content"
    assert details["provider_used"] == "mineru"
    assert details["parser_routing"]["mode"] == "auto"
    assert details["parser_routing"]["primary_provider"] == "mineru"
    assert details["provider_attempted"] == ["mineru"]
    assert calls == ["mineru"]


def test_extract_text_for_rag_auto_routes_word_to_llamaparse(
    tmp_path: Path, monkeypatch
):
    import services.file_parser as file_parser_module

    class _FakeLlamaProvider:
        name = "llamaparse"

        def supports(self, _file_type: str) -> bool:
            return True

        def extract_text(self, _filepath: str, _filename: str, _file_type: str):
            return "llamaparse content", {"pages_extracted": 1, "text_length": 17}

    calls: list[str] = []

    def _fake_get_parser(provider_name=None):
        calls.append(str(provider_name))
        if provider_name == "llamaparse":
            return _FakeLlamaProvider()
        raise AssertionError(f"unexpected provider request: {provider_name}")

    docx_path = tmp_path / "sample.docx"
    docx_path.write_bytes(b"fake-docx-content")

    monkeypatch.setenv("DOCUMENT_PARSER", "auto")
    monkeypatch.setattr(file_parser_module, "get_parser", _fake_get_parser)

    text, details = extract_text_for_rag(str(docx_path), "sample.docx", "word")

    assert text == "llamaparse content"
    assert details["provider_used"] == "llamaparse"
    assert details["parser_routing"]["mode"] == "auto"
    assert details["parser_routing"]["primary_provider"] == "llamaparse"
    assert details["provider_attempted"] == ["llamaparse"]
    assert calls == ["llamaparse"]
