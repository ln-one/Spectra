"""
LocalProvider — MVP 阶段的本地轻量解析 provider。

使用 pypdf / python-docx / python-pptx 在本地完成文档文本提取，
作为默认 provider（``DOCUMENT_PARSER=local``）。
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

from services.file_upload_service import FileType, normalize_file_type

from .base import BaseParseProvider


class LocalProvider(BaseParseProvider):
    """基于 pypdf / python-docx / python-pptx 的本地轻量解析器。"""

    name = "local"
    supported_types = {"pdf", "word", "ppt", "other"}

    def extract_text(
        self, filepath: str, filename: str, file_type: str
    ) -> tuple[str, dict[str, Any]]:
        path = Path(filepath)
        ext = path.suffix.lower()
        normalized_file_type = normalize_file_type(file_type)

        if normalized_file_type == FileType.PDF or ext == ".pdf":
            return self._extract_pdf(path)

        # 纯文本类型（txt/md/csv）优先按扩展名短路，避免错误走 _extract_docx
        if ext in {".txt", ".md", ".csv"}:
            text = path.read_text(encoding="utf-8", errors="replace")
            return text, {"text_length": len(text)}

        if normalized_file_type == FileType.WORD or ext in {".docx", ".doc"}:
            return self._extract_docx(path)

        if normalized_file_type == FileType.PPT or ext in {".pptx", ".ppt"}:
            return self._extract_pptx(path)

        # 未知类型：尝试按文本读取
        text = path.read_text(encoding="utf-8", errors="replace")
        return text, {"text_length": len(text)}

    # ------------------------------------------------------------------
    # 内部解析方法
    # ------------------------------------------------------------------

    @staticmethod
    def _extract_pdf(path: Path) -> tuple[str, dict[str, Any]]:
        details: dict[str, Any] = {"pages_extracted": 0, "text_length": 0}
        try:
            from pypdf import PdfReader

            reader = PdfReader(str(path))
            pages = [page.extract_text() or "" for page in reader.pages]
            text = "\n\n".join(pages).strip()
            details["pages_extracted"] = len(reader.pages)
            details["text_length"] = len(text)
            return text, details
        except Exception:
            # 解析失败返回空文本，由上层统一 fallback
            return "", details

    @staticmethod
    def _extract_docx(path: Path) -> tuple[str, dict[str, Any]]:
        details: dict[str, Any] = {"pages_extracted": 0, "text_length": 0}
        try:
            from docx import Document

            doc = Document(str(path))
            paragraphs = [
                p.text.strip() for p in doc.paragraphs if p.text and p.text.strip()
            ]
            text = "\n".join(paragraphs).strip()
            details["pages_extracted"] = max(1, text.count("\n\n"))
            details["text_length"] = len(text)
            return text, details
        except Exception:
            # 与 MVP 行为对齐：DOCX 解析失败时 pages_extracted 仍为 1
            details["pages_extracted"] = 1
            return "", details

    @staticmethod
    def _extract_pptx(path: Path) -> tuple[str, dict[str, Any]]:
        details: dict[str, Any] = {"pages_extracted": 0, "text_length": 0}
        try:
            from pptx import Presentation

            prs = Presentation(str(path))
            slide_texts: list[str] = []
            for slide in prs.slides:
                shape_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        content = shape.text.strip()
                        if content:
                            shape_texts.append(content)
                if shape_texts:
                    slide_texts.append("\n".join(shape_texts))
            text = "\n\n".join(slide_texts).strip()
            details["pages_extracted"] = len(prs.slides)
            details["text_length"] = len(text)
            return text, details
        except Exception:
            return "", details
