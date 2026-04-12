import logging
from pathlib import Path
from typing import Any, Dict
from uuid import uuid4

try:
    from docx import Document
except Exception:  # pragma: no cover - optional dependency guard
    Document = None

try:
    from pptx import Presentation
except Exception:  # pragma: no cover - optional dependency guard
    Presentation = None

from services.artifact_generator.office_placeholders import (
    generate_docx_placeholder,
    generate_pptx_placeholder,
)
from services.artifact_generator.policies import allow_office_placeholder_artifacts
from utils.docx_content_sidecar import write_docx_content_sidecar

logger = logging.getLogger(__name__)


class ArtifactOfficeMixin:
    @staticmethod
    def _should_emit_placeholder_artifact() -> bool:
        return allow_office_placeholder_artifacts()

    @staticmethod
    def _build_marp_markdown(content: Dict[str, Any], title: str) -> str:
        raw_markdown = str(content.get("markdown_content") or "").strip()
        if raw_markdown:
            return raw_markdown

        slides = content.get("slides")
        slide_items = slides if isinstance(slides, list) else []
        if not slide_items:
            slide_items = [{"title": title, "content": content.get("summary", "")}]

        blocks: list[str] = [
            "---",
            "marp: true",
            "theme: default",
            "paginate: true",
            "---",
        ]
        for item in slide_items:
            slide_title = str(item.get("title") or title).strip()
            slide_content = str(
                item.get("content")
                or item.get("description")
                or item.get("summary")
                or ""
            ).strip()
            blocks.append(f"# {slide_title or title}")
            if slide_content:
                blocks.append(slide_content)
            blocks.append("---")
        return "\n\n".join(blocks).rstrip("- \n")

    @staticmethod
    def _build_doc_markdown(content: Dict[str, Any], title: str) -> str:
        lesson_plan_markdown = str(content.get("lesson_plan_markdown") or "").strip()
        if lesson_plan_markdown:
            return lesson_plan_markdown

        lines: list[str] = [f"# {title}"]
        summary = str(content.get("summary") or "").strip()
        if summary:
            lines.extend(["", summary])

        sections = content.get("sections")
        section_items = sections if isinstance(sections, list) else []
        for section in section_items:
            section_title = str(section.get("title") or "").strip()
            section_content = str(
                section.get("content")
                or section.get("description")
                or section.get("summary")
                or ""
            ).strip()
            if section_title:
                lines.extend(["", f"## {section_title}"])
            if section_content:
                lines.extend(["", section_content])
        return "\n".join(lines).strip()

    async def _render_pptx_with_marp(self, storage_path: str, markdown: str) -> bool:
        try:
            from services.generation.marp_generator import call_marp_cli
        except Exception as exc:
            logger.debug("Marp renderer unavailable in artifact generator: %s", exc)
            return False

        output_path = Path(storage_path)
        temp_md = output_path.with_name(f".tmp-{uuid4().hex}.md")
        temp_md.write_text(markdown, encoding="utf-8")
        try:
            await call_marp_cli(temp_md, output_path)
        except Exception as exc:
            logger.warning("Marp render failed, fallback to local generator: %s", exc)
            return False
        finally:
            temp_md.unlink(missing_ok=True)
        return output_path.exists() and output_path.stat().st_size > 0

    async def _render_docx_with_pandoc(self, storage_path: str, markdown: str) -> bool:
        try:
            from services.generation.pandoc_generator import call_pandoc
        except Exception as exc:
            logger.debug("Pandoc renderer unavailable in artifact generator: %s", exc)
            return False

        output_path = Path(storage_path)
        temp_md = output_path.with_name(f".tmp-{uuid4().hex}.md")
        temp_md.write_text(markdown, encoding="utf-8")
        try:
            await call_pandoc(temp_md, output_path)
        except Exception as exc:
            logger.warning("Pandoc render failed, fallback to local generator: %s", exc)
            return False
        finally:
            temp_md.unlink(missing_ok=True)
        return output_path.exists() and output_path.stat().st_size > 0

    async def _render_docx_with_pandoc_html(self, storage_path: str, html_content: str) -> bool:
        try:
            from services.generation.pandoc_generator import call_pandoc
        except Exception as exc:
            logger.debug("Pandoc renderer unavailable in artifact generator: %s", exc)
            return False

        output_path = Path(storage_path)
        temp_html = output_path.with_name(f".tmp-{uuid4().hex}.html")
        temp_html.write_text(html_content, encoding="utf-8")
        try:
            await call_pandoc(temp_html, output_path)
        except Exception as exc:
            logger.warning("Pandoc HTML render failed, fallback to local generator: %s", exc)
            return False
        finally:
            temp_html.unlink(missing_ok=True)
        return output_path.exists() and output_path.stat().st_size > 0

    async def generate_pptx(
        self, content: Dict[str, Any], project_id: str, artifact_id: str
    ) -> str:
        storage_path = self.get_storage_path(project_id, "pptx", artifact_id)
        title = str(content.get("title", "Project Space PPTX")).strip()
        markdown = self._build_marp_markdown(content, title)

        if await self._render_pptx_with_marp(storage_path, markdown):
            logger.info("Generated PPTX artifact via Marp at %s", storage_path)
            return storage_path

        if Presentation is not None:
            try:
                presentation = Presentation()
                slides = content.get("slides")
                slide_items = slides if isinstance(slides, list) else []
                if not slide_items:
                    slide_items = [{"title": title, "content": ""}]

                for item in slide_items:
                    slide_title = str(item.get("title") or title).strip()
                    slide_content = str(
                        item.get("content")
                        or item.get("description")
                        or item.get("summary")
                        or ""
                    ).strip()
                    layout = presentation.slide_layouts[1]
                    slide = presentation.slides.add_slide(layout)
                    slide.shapes.title.text = slide_title
                    placeholder = slide.placeholders[1]
                    placeholder.text = slide_content

                presentation.save(storage_path)
                logger.info("Generated PPTX artifact at %s", storage_path)
                return storage_path
            except Exception as exc:
                logger.warning(
                    "python-pptx generation failed, fallback to placeholder: %s", exc
                )
        if self._should_emit_placeholder_artifact():
            generate_pptx_placeholder(storage_path, title)
            logger.warning(
                (
                    "Generated PPTX placeholder artifact at %s "
                    "because no renderer succeeded"
                ),
                storage_path,
            )
            return storage_path
        raise RuntimeError(
            "PPTX rendering failed and placeholder artifacts are disabled. "
            "Enable ALLOW_OFFICE_PLACEHOLDER_ARTIFACTS only for explicit dev fallback."
        )

    async def generate_docx(
        self, content: Dict[str, Any], project_id: str, artifact_id: str
    ) -> str:
        storage_path = self.get_storage_path(project_id, "docx", artifact_id)
        title = str(content.get("title", "Project Space DOCX")).strip()
        markdown = self._build_doc_markdown(content, title)
        doc_source_html = str(content.get("doc_source_html") or "").strip()

        if doc_source_html and await self._render_docx_with_pandoc_html(
            storage_path, doc_source_html
        ):
            write_docx_content_sidecar(storage_path, content)
            logger.info("Generated DOCX artifact via Pandoc HTML at %s", storage_path)
            return storage_path

        if await self._render_docx_with_pandoc(storage_path, markdown):
            write_docx_content_sidecar(storage_path, content)
            logger.info("Generated DOCX artifact via Pandoc at %s", storage_path)
            return storage_path

        if Document is not None:
            try:
                document = Document()
                document.add_heading(title, level=1)
                sections = content.get("sections")
                section_items = sections if isinstance(sections, list) else []
                if not section_items and content.get("summary"):
                    section_items = [
                        {"title": "Summary", "content": content["summary"]}
                    ]

                for item in section_items:
                    section_title = str(item.get("title") or "").strip()
                    section_content = str(
                        item.get("content")
                        or item.get("description")
                        or item.get("summary")
                        or ""
                    ).strip()
                    if section_title:
                        document.add_heading(section_title, level=2)
                    if section_content:
                        document.add_paragraph(section_content)

                document.save(storage_path)
                write_docx_content_sidecar(storage_path, content)
                logger.info("Generated DOCX artifact at %s", storage_path)
                return storage_path
            except Exception as exc:
                logger.warning(
                    "python-docx generation failed, fallback to placeholder: %s", exc
                )
        if self._should_emit_placeholder_artifact():
            generate_docx_placeholder(storage_path, title)
            logger.warning(
                (
                    "Generated DOCX placeholder artifact at %s "
                    "because no renderer succeeded"
                ),
                storage_path,
            )
            return storage_path
        raise RuntimeError(
            "DOCX rendering failed and placeholder artifacts are disabled. "
            "Enable ALLOW_OFFICE_PLACEHOLDER_ARTIFACTS only for explicit dev fallback."
        )
