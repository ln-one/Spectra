from pathlib import Path
from types import SimpleNamespace

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from services.generation_session_service.ppt_animation_binding import (
    apply_animation_placement_to_ppt_artifact,
)


def _build_gif(path: Path) -> None:
    frame_1 = Image.new("RGB", (320, 180), "#1f7a8c")
    frame_2 = Image.new("RGB", (320, 180), "#bfdbf7")
    frame_1.save(
        path,
        format="GIF",
        save_all=True,
        append_images=[frame_2],
        duration=160,
        loop=0,
    )


def _build_pptx(path: Path, slide_count: int = 3) -> None:
    presentation = Presentation()
    for index in range(slide_count):
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = f"Slide {index + 1}"
        slide.placeholders[1].text = f"Content {index + 1}"
    presentation.save(str(path))


def _picture_shapes(slide) -> list:
    return [
        shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    ]


def test_apply_animation_placement_to_ppt_artifact_inserts_gif_on_target_slides(
    tmp_path: Path,
):
    gif_path = tmp_path / "demo.gif"
    ppt_path = tmp_path / "deck.pptx"
    _build_gif(gif_path)
    _build_pptx(ppt_path, slide_count=3)

    animation_artifact = SimpleNamespace(id="anim-001", storagePath=str(gif_path))
    ppt_artifact = SimpleNamespace(id="ppt-001", storagePath=str(ppt_path))

    apply_animation_placement_to_ppt_artifact(
        animation_artifact=animation_artifact,
        ppt_artifact=ppt_artifact,
        placement_records=[
            {"page_number": 2, "slot": "bottom-right"},
            {"page_number": 3, "slot": "right-panel"},
        ],
    )

    presentation = Presentation(str(ppt_path))
    assert len(_picture_shapes(presentation.slides[0])) == 0
    assert len(_picture_shapes(presentation.slides[1])) == 1
    assert len(_picture_shapes(presentation.slides[2])) == 1


def test_apply_animation_placement_to_ppt_artifact_replaces_existing_picture_on_same_slide(
    tmp_path: Path,
):
    gif_path = tmp_path / "demo.gif"
    ppt_path = tmp_path / "deck.pptx"
    _build_gif(gif_path)
    _build_pptx(ppt_path, slide_count=2)

    animation_artifact = SimpleNamespace(id="anim-001", storagePath=str(gif_path))
    ppt_artifact = SimpleNamespace(id="ppt-001", storagePath=str(ppt_path))

    apply_animation_placement_to_ppt_artifact(
        animation_artifact=animation_artifact,
        ppt_artifact=ppt_artifact,
        placement_records=[{"page_number": 2, "slot": "bottom-right"}],
    )
    initial_presentation = Presentation(str(ppt_path))
    initial_picture = _picture_shapes(initial_presentation.slides[1])[0]
    initial_left = initial_picture.left

    apply_animation_placement_to_ppt_artifact(
        animation_artifact=animation_artifact,
        ppt_artifact=ppt_artifact,
        placement_records=[{"page_number": 2, "slot": "right-panel"}],
    )

    updated_presentation = Presentation(str(ppt_path))
    pictures = _picture_shapes(updated_presentation.slides[1])
    assert len(pictures) == 1
    assert pictures[0].left != initial_left
