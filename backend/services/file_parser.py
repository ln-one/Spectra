"""
File Parser Service

MVP 阶段的轻量解析服务：
- PDF: pypdf
- DOCX: python-docx
- PPTX: python-pptx
- 纯文本: 直接读取
- 图片/视频: 生成可检索的占位描述
"""

from __future__ import annotations

from pathlib import Path
from typing import Any


def extract_text_for_rag(
    filepath: str, filename: str, file_type: str
) -> tuple[str, dict]:
    """
    从文件中提取可用于 RAG 的文本及解析详情。

    Returns:
        (text, parse_details)
    """
    path = Path(filepath)
    ext = path.suffix.lower()
    details: dict[str, Any] = {}

    if file_type == "pdf" or ext == ".pdf":
        text, pdf_details = _extract_pdf_text(path)
        details.update(pdf_details)
        return text, details

    if ext in {".txt", ".md", ".csv"}:
        text = path.read_text(encoding="utf-8", errors="replace")
        details["text_length"] = len(text)
        return text, details

    if file_type == "word" or ext in {".docx", ".doc"}:
        text = _extract_docx_text(path)
        details["pages_extracted"] = max(1, text.count("\n\n"))
        details["text_length"] = len(text)
        return text, details

    if file_type == "ppt" or ext in {".pptx", ".ppt"}:
        text, slide_count = _extract_pptx_text(path)
        details["pages_extracted"] = slide_count
        details["text_length"] = len(text)
        return text, details

    # 多模态占位信息：至少将文件元信息入库，便于被检索和引用
    if file_type == "image":
        text = f"图片资料：{filename}。该图片可作为课堂讲解示例或视觉辅助素材。"
        details["images_extracted"] = 1
        details["text_length"] = len(text)
        return text, details

    if file_type == "video":
        text = f"视频资料：{filename}。该视频可用于课堂案例演示与讨论。"
        details["duration"] = 0
        details["text_length"] = len(text)
        return text, details

    # 兜底：尝试按文本读取
    text = path.read_text(encoding="utf-8", errors="replace")
    details["text_length"] = len(text)
    return text, details


def _extract_pdf_text(path: Path) -> tuple[str, dict]:
    details: dict[str, Any] = {"pages_extracted": 0, "text_length": 0}
    try:
        from pypdf import PdfReader

        reader = PdfReader(str(path))
        pages = []
        for page in reader.pages:
            pages.append(page.extract_text() or "")
        text = "\n\n".join(pages).strip()
        details["pages_extracted"] = len(reader.pages)
        details["text_length"] = len(text)
        return text, details
    except Exception:
        # 失败时返回空文本，由上层决定是否降级占位
        return "", details


def _extract_docx_text(path: Path) -> str:
    try:
        from docx import Document

        doc = Document(str(path))
        paragraphs = [
            p.text.strip() for p in doc.paragraphs if p.text and p.text.strip()
        ]
        return "\n".join(paragraphs).strip()
    except Exception:
        return ""


def _extract_pptx_text(path: Path) -> tuple[str, int]:
    try:
        from pptx import Presentation

        prs = Presentation(str(path))
        slide_texts: list[str] = []
        for slide in prs.slides:
            shape_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    content = shape.text.strip()
                    if content:
                        shape_texts.append(content)
            if shape_texts:
                slide_texts.append("\n".join(shape_texts))
        return "\n\n".join(slide_texts).strip(), len(prs.slides)
    except Exception:
        return "", 0
