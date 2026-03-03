"""
file_parser.py 回归测试。

验证 extract_text_for_rag() 在可插拔架构改造后行为与 MVP 完全一致。
测试文件在 tmp_path 中动态生成，不提交二进制 fixture。
"""

from pathlib import Path

from services.file_parser import extract_text_for_rag

# ---- 纯文本 ----


def test_extract_text_for_rag_txt_file(tmp_path: Path):
    file_path = tmp_path / "notes.txt"
    file_path.write_text("第一行内容\n第二行内容", encoding="utf-8")

    text, details = extract_text_for_rag(str(file_path), "notes.txt", "other")

    assert "第一行内容" in text
    assert details["text_length"] > 0


def test_extract_text_for_rag_md_file(tmp_path: Path):
    file_path = tmp_path / "readme.md"
    file_path.write_text("# 标题\n正文内容", encoding="utf-8")

    text, details = extract_text_for_rag(str(file_path), "readme.md", "other")

    assert "标题" in text
    assert details["text_length"] > 0


def test_extract_text_for_rag_csv_file(tmp_path: Path):
    file_path = tmp_path / "data.csv"
    file_path.write_text("姓名,分数\n张三,90", encoding="utf-8")

    text, details = extract_text_for_rag(str(file_path), "data.csv", "other")

    assert "张三" in text
    assert details["text_length"] > 0


# ---- 图片/视频占位 ----


def test_extract_text_for_rag_image_placeholder(tmp_path: Path):
    file_path = tmp_path / "figure.png"
    file_path.write_bytes(b"\x89PNG\r\n")

    text, details = extract_text_for_rag(str(file_path), "figure.png", "image")

    assert "图片资料" in text
    assert "figure.png" in text
    assert details["images_extracted"] == 1


def test_extract_text_for_rag_video_placeholder(tmp_path: Path):
    file_path = tmp_path / "lecture.mp4"
    file_path.write_bytes(b"\x00\x00\x00\x1cftyp")

    text, details = extract_text_for_rag(str(file_path), "lecture.mp4", "video")

    assert "视频资料" in text
    assert "lecture.mp4" in text
    assert details["duration"] == 0
    assert details["text_length"] > 0


# ---- PDF 动态生成 ----


def test_extract_text_for_rag_pdf(tmp_path: Path):
    """动态生成最小 PDF 并验证解析。"""
    from pypdf import PdfWriter

    writer = PdfWriter()
    writer.add_blank_page(width=72, height=72)
    pdf_path = tmp_path / "test.pdf"
    with open(pdf_path, "wb") as f:
        writer.write(f)

    text, details = extract_text_for_rag(str(pdf_path), "test.pdf", "pdf")

    assert isinstance(text, str)
    assert details["pages_extracted"] == 1
    assert "text_length" in details


def test_extract_text_for_rag_pdf_corrupted(tmp_path: Path):
    """损坏 PDF 应返回空文本，不抛异常。"""
    pdf_path = tmp_path / "bad.pdf"
    pdf_path.write_bytes(b"not a pdf content at all")

    text, details = extract_text_for_rag(str(pdf_path), "bad.pdf", "pdf")

    assert text == ""
    assert details["pages_extracted"] == 0


# ---- DOCX 动态生成 ----


def test_extract_text_for_rag_docx(tmp_path: Path):
    """动态生成最小 DOCX 并验证解析。"""
    from docx import Document

    doc = Document()
    doc.add_paragraph("测试段落一")
    doc.add_paragraph("测试段落二")
    docx_path = tmp_path / "test.docx"
    doc.save(str(docx_path))

    text, details = extract_text_for_rag(str(docx_path), "test.docx", "word")

    assert "测试段落一" in text
    assert "测试段落二" in text
    assert details["pages_extracted"] >= 1
    assert details["text_length"] > 0


def test_extract_text_for_rag_docx_corrupted(tmp_path: Path):
    """损坏 DOCX 应返回空文本，不抛异常。"""
    docx_path = tmp_path / "bad.docx"
    docx_path.write_bytes(b"not a docx")

    text, details = extract_text_for_rag(str(docx_path), "bad.docx", "word")

    assert text == ""
    # 与 MVP 行为一致：DOCX 失败时仍返回 pages_extracted=1
    assert details["pages_extracted"] == 1


# ---- PPTX 动态生成 ----


def test_extract_text_for_rag_pptx(tmp_path: Path):
    """动态生成最小 PPTX 并验证解析。"""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    txBox.text_frame.text = "幻灯片文本内容"
    pptx_path = tmp_path / "test.pptx"
    prs.save(str(pptx_path))

    text, details = extract_text_for_rag(str(pptx_path), "test.pptx", "ppt")

    assert "幻灯片文本内容" in text
    assert details["pages_extracted"] == 1
    assert details["text_length"] > 0


def test_extract_text_for_rag_pptx_corrupted(tmp_path: Path):
    """损坏 PPTX 应返回空文本，不抛异常。"""
    pptx_path = tmp_path / "bad.pptx"
    pptx_path.write_bytes(b"not a pptx")

    text, details = extract_text_for_rag(str(pptx_path), "bad.pptx", "ppt")

    assert text == ""
    assert details["pages_extracted"] == 0


# ---- 接口签名兼容性 ----


def test_extract_text_for_rag_signature():
    """确保函数签名与 MVP 一致。"""
    import inspect

    sig = inspect.signature(extract_text_for_rag)
    params = list(sig.parameters.keys())
    assert params == ["filepath", "filename", "file_type"]


def test_extract_text_for_rag_fallback_to_local_when_provider_unsupported(
    tmp_path: Path, monkeypatch
):
    """当配置 provider 不支持 file_type 时，应自动回退到 local。"""
    import services.file_parser as file_parser_module
    from services.parsers.local_provider import LocalProvider

    class _UnsupportedProvider:
        name = "mineru"

        def supports(self, file_type: str) -> bool:
            return False

        def extract_text(self, filepath: str, filename: str, file_type: str):
            raise AssertionError("unsupported provider should not be used")

    calls = {"local": 0}

    def _fake_get_parser(provider_name=None):
        if provider_name == "local":
            calls["local"] += 1
            return LocalProvider()
        return _UnsupportedProvider()

    monkeypatch.setattr(file_parser_module, "get_parser", _fake_get_parser)

    file_path = tmp_path / "notes.txt"
    file_path.write_text("fallback local content", encoding="utf-8")

    text, details = extract_text_for_rag(str(file_path), "notes.txt", "other")

    assert "fallback local content" in text
    assert details["text_length"] > 0
    assert calls["local"] == 1
