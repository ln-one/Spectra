"""
Parser providers 单元测试。

覆盖 provider 注册/切换、fallback 机制、BaseParseProvider 约定等。
使用 monkeypatch 模拟 provider 不可用 / 环境变量切换等场景。
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

from services.parsers.base import BaseParseProvider, ProviderNotAvailableError
from services.parsers.local_provider import LocalProvider
from services.parsers.registry import _PROVIDER_FACTORIES, get_parser, register_provider

# =====================================================================
# BaseParseProvider 接口约定
# =====================================================================


class _DummyProvider(BaseParseProvider):
    name = "dummy"
    supported_types = {"pdf", "word"}

    def extract_text(
        self, filepath: str, filename: str, file_type: str
    ) -> tuple[str, dict[str, Any]]:
        return "dummy text", {"text_length": 10}


class TestBaseParseProvider:
    def test_supports_returns_true_for_supported_type(self):
        p = _DummyProvider()
        assert p.supports("pdf") is True
        assert p.supports("word") is True

    def test_supports_returns_false_for_unsupported_type(self):
        p = _DummyProvider()
        assert p.supports("image") is False
        assert p.supports("video") is False

    def test_extract_text_signature(self):
        p = _DummyProvider()
        text, details = p.extract_text("/fake/path", "file.pdf", "pdf")
        assert isinstance(text, str)
        assert isinstance(details, dict)


# =====================================================================
# LocalProvider
# =====================================================================


class TestLocalProvider:
    def test_name_and_supported_types(self):
        p = LocalProvider()
        assert p.name == "local"
        assert "pdf" in p.supported_types
        assert "word" in p.supported_types
        assert "ppt" in p.supported_types

    def test_pdf_extraction(self, tmp_path: Path):
        from pypdf import PdfWriter

        writer = PdfWriter()
        writer.add_blank_page(72, 72)
        pdf_path = tmp_path / "test.pdf"
        with open(pdf_path, "wb") as f:
            writer.write(f)

        p = LocalProvider()
        text, details = p.extract_text(str(pdf_path), "test.pdf", "pdf")
        assert isinstance(text, str)
        assert details["pages_extracted"] == 1

    def test_docx_extraction(self, tmp_path: Path):
        from docx import Document

        doc = Document()
        doc.add_paragraph("本地解析测试")
        docx_path = tmp_path / "test.docx"
        doc.save(str(docx_path))

        p = LocalProvider()
        text, details = p.extract_text(str(docx_path), "test.docx", "word")
        assert "本地解析测试" in text
        assert details["pages_extracted"] >= 1

    def test_pptx_extraction(self, tmp_path: Path):
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
        txBox.text_frame.text = "Provider 测试"
        pptx_path = tmp_path / "test.pptx"
        prs.save(str(pptx_path))

        p = LocalProvider()
        text, details = p.extract_text(str(pptx_path), "test.pptx", "ppt")
        assert "Provider 测试" in text
        assert details["pages_extracted"] == 1

    def test_txt_fallback(self, tmp_path: Path):
        txt_path = tmp_path / "notes.txt"
        txt_path.write_text("纯文本内容", encoding="utf-8")

        p = LocalProvider()
        text, details = p.extract_text(str(txt_path), "notes.txt", "other")
        assert "纯文本内容" in text
        assert details["text_length"] > 0

    def test_corrupted_pdf(self, tmp_path: Path):
        bad_path = tmp_path / "bad.pdf"
        bad_path.write_bytes(b"corrupt")

        p = LocalProvider()
        text, details = p.extract_text(str(bad_path), "bad.pdf", "pdf")
        assert text == ""
        assert details["pages_extracted"] == 0

    def test_corrupted_docx(self, tmp_path: Path):
        bad_path = tmp_path / "bad.docx"
        bad_path.write_bytes(b"corrupt")

        p = LocalProvider()
        text, details = p.extract_text(str(bad_path), "bad.docx", "word")
        assert text == ""
        assert details["pages_extracted"] == 1

    def test_corrupted_pptx(self, tmp_path: Path):
        bad_path = tmp_path / "bad.pptx"
        bad_path.write_bytes(b"corrupt")

        p = LocalProvider()
        text, details = p.extract_text(str(bad_path), "bad.pptx", "ppt")
        assert text == ""
        assert details["pages_extracted"] == 0


# =====================================================================
# Registry: get_parser / 切换 / fallback
# =====================================================================


class TestRegistry:
    def test_default_is_local(self, monkeypatch):
        monkeypatch.delenv("DOCUMENT_PARSER", raising=False)
        parser = get_parser()
        assert parser.name == "local"

    def test_explicit_local(self):
        parser = get_parser("local")
        assert parser.name == "local"

    def test_explicit_auto_uses_local_fallback_provider(self):
        parser = get_parser("auto")
        assert parser.name == "local"

    def test_env_var_selects_provider(self, monkeypatch):
        monkeypatch.setenv("DOCUMENT_PARSER", "local")
        parser = get_parser()
        assert parser.name == "local"

    def test_unknown_provider_falls_back_to_local(self, monkeypatch):
        monkeypatch.setenv("DOCUMENT_PARSER", "nonexistent")
        parser = get_parser()
        assert parser.name == "local"

    def test_unavailable_provider_falls_back_to_local(self, monkeypatch):
        """模拟 mineru 未安装的场景。"""

        def _raise_unavailable():
            raise ProviderNotAvailableError("mocked")

        monkeypatch.setitem(_PROVIDER_FACTORIES, "mineru", _raise_unavailable)
        parser = get_parser("mineru")
        assert parser.name == "local"

    def test_register_custom_provider(self):
        """测试自定义 provider 注册。"""
        register_provider("dummy", lambda: _DummyProvider())
        try:
            parser = get_parser("dummy")
            assert parser.name == "dummy"
            text, details = parser.extract_text("", "", "pdf")
            assert text == "dummy text"
        finally:
            # 清理注册
            _PROVIDER_FACTORIES.pop("dummy", None)


# =====================================================================
# parse_details 标准字段验证（与 OpenAPI 对齐）
# =====================================================================

# OpenAPI files.yaml 中 parse_details 的标准字段集合
_STANDARD_FIELDS = {"pages_extracted", "images_extracted", "text_length", "duration"}


class TestParseDetailsContract:
    """确保 LocalProvider 返回的 details 字段与 OpenAPI parse_details 契约一致。"""

    def test_pdf_details_fields(self, tmp_path: Path):
        from pypdf import PdfWriter

        writer = PdfWriter()
        writer.add_blank_page(72, 72)
        pdf_path = tmp_path / "test.pdf"
        with open(pdf_path, "wb") as f:
            writer.write(f)

        p = LocalProvider()
        _, details = p.extract_text(str(pdf_path), "test.pdf", "pdf")
        assert "pages_extracted" in details
        assert "text_length" in details
        for key in details:
            assert key in _STANDARD_FIELDS, f"非标准字段: {key}"

    def test_docx_details_fields(self, tmp_path: Path):
        from docx import Document

        doc = Document()
        doc.add_paragraph("字段测试")
        docx_path = tmp_path / "test.docx"
        doc.save(str(docx_path))

        p = LocalProvider()
        _, details = p.extract_text(str(docx_path), "test.docx", "word")
        assert "pages_extracted" in details
        assert "text_length" in details
        for key in details:
            assert key in _STANDARD_FIELDS, f"非标准字段: {key}"

    def test_pptx_details_fields(self, tmp_path: Path):
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
        txBox.text_frame.text = "字段测试"
        pptx_path = tmp_path / "test.pptx"
        prs.save(str(pptx_path))

        p = LocalProvider()
        _, details = p.extract_text(str(pptx_path), "test.pptx", "ppt")
        assert "pages_extracted" in details
        assert "text_length" in details
        for key in details:
            assert key in _STANDARD_FIELDS, f"非标准字段: {key}"
